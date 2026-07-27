"""
Unified Multiformat Parser for RAG Ingestion.
Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT, PNG, JPEG, WEBP.
Extracted text is mapped to page-like structures with metadata.
"""

from __future__ import annotations

import logging
import os
import hashlib
import csv
from typing import Any

logger = logging.getLogger(__name__)

# Fallback parser if external libraries are missing
def compute_file_hash(file_path: str) -> str:
    """Compute SHA-256 hash of a file to detect duplicates."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        while chunk := f.read(8192):
            sha256.update(chunk)
    return sha256.hexdigest()

def extract_txt(file_path: str) -> dict[str, Any]:
    """Parse plain text files."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return {
            "full_text": text,
            "pages": [{
                "page_number": 1,
                "text": text,
                "pdf_text": text,
                "ocr_text": "",
                "ocr_used": False,
                "images": []
            }],
            "total_pages": 1
        }
    except Exception as exc:
        logger.error("Failed to parse TXT %s: %s", file_path, exc)
        raise

def extract_csv(file_path: str) -> dict[str, Any]:
    """Parse CSV files and format rows as markdown tables."""
    try:
        markdown_rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = list(reader)
        
        if not rows:
            return {"full_text": "", "pages": [], "total_pages": 0}
            
        # Format as Markdown table
        headers = rows[0]
        markdown_rows.append("| " + " | ".join(headers) + " |")
        markdown_rows.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in rows[1:]:
            # Pad row if columns don't match
            if len(row) < len(headers):
                row = row + [""] * (len(headers) - len(row))
            elif len(row) > len(headers):
                row = row[:len(headers)]
            markdown_rows.append("| " + " | ".join(row) + " |")
            
        table_text = "\n".join(markdown_rows)
        return {
            "full_text": table_text,
            "pages": [{
                "page_number": 1,
                "text": table_text,
                "pdf_text": table_text,
                "ocr_text": "",
                "ocr_used": False,
                "images": []
            }],
            "total_pages": 1
        }
    except Exception as exc:
        logger.error("Failed to parse CSV %s: %s", file_path, exc)
        raise

def extract_docx(file_path: str) -> dict[str, Any]:
    """Parse Word documents (DOCX)."""
    try:
        import docx
    except ImportError:
        logger.warning("python-docx is not installed. Falling back to basic text decoding.")
        return extract_txt(file_path)
        
    try:
        doc = docx.Document(file_path)
        full_text_parts = []
        
        # Read paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                full_text_parts.append(para.text)
                
        # Read tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                if any(row_text):
                    full_text_parts.append(" | ".join(row_text))
                    
        full_text = "\n".join(full_text_parts)
        
        # Docx doesn't have strict page boundaries natively, we split by ~500 words to create virtual pages
        words = full_text.split()
        pages = []
        chunk_size = 500
        for i in range(0, len(words), chunk_size):
            page_words = words[i:i+chunk_size]
            page_text = " ".join(page_words)
            page_num = (i // chunk_size) + 1
            pages.append({
                "page_number": page_num,
                "text": page_text,
                "pdf_text": page_text,
                "ocr_text": "",
                "ocr_used": False,
                "images": []
            })
            
        return {
            "full_text": full_text,
            "pages": pages,
            "total_pages": len(pages)
        }
    except Exception as exc:
        logger.error("Failed to parse DOCX %s: %s", file_path, exc)
        raise

def extract_pptx(file_path: str) -> dict[str, Any]:
    """Parse PowerPoint presentations (PPTX)."""
    try:
        import pptx
    except ImportError:
        logger.warning("python-pptx is not installed. Falling back to basic text decoding.")
        return extract_txt(file_path)
        
    try:
        prs = pptx.Presentation(file_path)
        pages = []
        full_text_parts = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            slide_text = "\n".join(slide_text_parts)
            page_number = slide_idx + 1
            full_text_parts.append(f"\n--- Slide {page_number} ---\n{slide_text}")
            
            pages.append({
                "page_number": page_number,
                "text": slide_text,
                "pdf_text": slide_text,
                "ocr_text": "",
                "ocr_used": False,
                "images": []
            })
            
        return {
            "full_text": "\n".join(full_text_parts),
            "pages": pages,
            "total_pages": len(pages)
        }
    except Exception as exc:
        logger.error("Failed to parse PPTX %s: %s", file_path, exc)
        raise

def extract_xlsx(file_path: str) -> dict[str, Any]:
    """Parse Excel spreadsheets (XLSX)."""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl is not installed. Falling back to basic text decoding.")
        return extract_txt(file_path)
        
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        pages = []
        full_text_parts = []
        page_counter = 1
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                row_str = [str(cell).strip() if cell is not None else "" for cell in row]
                if any(row_str):
                    sheet_rows.append("| " + " | ".join(row_str) + " |")
            
            if sheet_rows:
                sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(sheet_rows)
                full_text_parts.append(sheet_text)
                pages.append({
                    "page_number": page_counter,
                    "text": sheet_text,
                    "pdf_text": sheet_text,
                    "ocr_text": "",
                    "ocr_used": False,
                    "images": []
                })
                page_counter += 1
                
        return {
            "full_text": "\n\n".join(full_text_parts),
            "pages": pages,
            "total_pages": len(pages)
        }
    except Exception as exc:
        logger.error("Failed to parse XLSX %s: %s", file_path, exc)
        raise

def extract_image(file_path: str, document_id: str, image_base_path: str) -> dict[str, Any]:
    """Parse image files (PNG, JPEG, WEBP) directly using PaddleOCR."""
    from ingestion.ocr_service import run_ocr_on_image
    
    try:
        # Run OCR on the image file
        ocr_result = run_ocr_on_image(file_path)
        text = ocr_result.text
        confidence = ocr_result.confidence
        
        image_id = "page_1"
        image_record = {
            "image_id": image_id,
            "page_number": 1,
            "page": 1,
            "image_path": f"{image_base_path}/{image_id}",
            "local_path": file_path,
        }
        
        return {
            "full_text": text,
            "pages": [{
                "page_number": 1,
                "text": text,
                "pdf_text": "",
                "ocr_text": text,
                "ocr_confidence": confidence,
                "ocr_used": True,
                "images": [image_record]
            }],
            "total_pages": 1
        }
    except Exception as exc:
        logger.error("Failed to parse image %s: %s", file_path, exc)
        raise

def parse_document(
    file_path: str,
    document_id: str,
    source_file: str,
    output_folder: str | None = None,
    image_base_path: str | None = None,
) -> dict[str, Any]:
    """
    Unified entry point for parsing any supported document file.
    Delegates to format-specific parser based on file extension.
    """
    ext = os.path.splitext(source_file)[1].lower()
    
    if image_base_path is None:
        image_base_path = f"documents/{document_id}/images"
        
    if ext == ".pdf":
        from ingestion.pdf_processor import process_pdf
        return process_pdf(
            pdf_path=file_path,
            document_id=document_id,
            source_file=source_file,
            output_folder=output_folder,
            image_base_path=image_base_path,
        )
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        return extract_image(file_path, document_id, image_base_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".xlsx":
        return extract_xlsx(file_path)
    elif ext == ".csv":
        return extract_csv(file_path)
    else:
        # Fallback to plain text
        return extract_txt(file_path)
